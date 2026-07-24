"""Convert a Compete Deck outline (Markdown) to a PowerPoint deck.

Slide-markdown convention: slides are separated by a line containing only
`---`. Each slide's first `#`/`##` line becomes the slide title; everything
else becomes body content. Handles bullets (with one level of nesting via
2-space indent), pipe tables, and bold spans. The very first slide is
rendered as a title slide (title + subtitle, no bullets). Not a general
Markdown-to-slides renderer -- this repo's deck format is deliberately
constrained, same spirit as md_to_docx.py.
"""
import pathlib
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TITLE_LAYOUT = 0
TITLE_CONTENT_LAYOUT = 1

DARK = RGBColor(0x23, 0x27, 0x2E)
ACCENT = RGBColor(0x33, 0x66, 0xCC)


def add_bold_runs(paragraph, text):
    for chunk in re.split(r"(\*\*.+?\*\*)", text):
        if chunk.startswith("**") and chunk.endswith("**"):
            run = paragraph.add_run()
            run.text = chunk[2:-2]
            run.font.bold = True
        elif chunk:
            run = paragraph.add_run()
            run.text = chunk


def is_table_row(line: str) -> bool:
    return line.strip().startswith("|") and line.strip().endswith("|")


def is_table_sep(line: str) -> bool:
    return bool(re.match(r"^\s*\|?[\s:|-]+\|?\s*$", line)) and "-" in line


def parse_slides(md_text: str):
    blocks = re.split(r"\n-{3,}\n", md_text.strip())
    slides = []
    for block in blocks:
        lines = [l for l in block.splitlines()]
        title = ""
        subtitle = ""
        body_lines = []
        i = 0
        while i < len(lines) and not lines[i].strip():
            i += 1
        if i < len(lines):
            m = re.match(r"^#{1,2}\s+(.*)", lines[i].strip())
            if m:
                title = m.group(1)
                i += 1
        if i < len(lines) and lines[i].strip().startswith("_") and lines[i].strip().endswith("_") and len(lines[i].strip()) > 1:
            subtitle = lines[i].strip().strip("_")
            i += 1
        body_lines = lines[i:]
        slides.append((title, subtitle, body_lines))
    return slides


def add_table(slide, table_lines, top):
    rows = [
        [c.strip() for c in row.strip("|").split("|")]
        for row in table_lines
        if not is_table_sep(row)
    ]
    if not rows:
        return top
    nrows, ncols = len(rows), len(rows[0])
    height = Inches(0.4 * nrows)
    table_shape = slide.shapes.add_table(
        nrows, ncols, Inches(0.5), top, Inches(9), height
    )
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx in range(ncols):
            text = row[c_idx] if c_idx < len(row) else ""
            cell = table.cell(r_idx, c_idx)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            add_bold_runs(p, text)
            for run in p.runs:
                run.font.size = Pt(12)
                if r_idx == 0:
                    run.font.bold = True
    return top + height + Inches(0.2)


def build_body(slide, body_lines):
    placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            placeholder = shape
            break
    tf = placeholder.text_frame if placeholder is not None else None
    if tf is not None:
        tf.clear()
    first_para = True

    i = 0
    extra_top = Inches(1.5)
    while i < len(body_lines):
        line = body_lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue

        if is_table_row(stripped):
            table_lines = []
            while i < len(body_lines) and is_table_row(body_lines[i].strip()):
                table_lines.append(body_lines[i].strip())
                i += 1
            extra_top = add_table(slide, table_lines, extra_top)
            continue

        bullet_m = re.match(r"^(\s*)-\s+(.*)", line)
        if bullet_m and tf is not None:
            indent = len(bullet_m.group(1))
            level = 1 if indent >= 2 else 0
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.level = level
            add_bold_runs(p, bullet_m.group(2))
            for run in p.runs:
                run.font.size = Pt(16 if level == 0 else 14)
            i += 1
            continue

        if tf is not None:
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            add_bold_runs(p, stripped)
            for run in p.runs:
                run.font.size = Pt(16)
            i += 1
            continue
        i += 1


def convert(md_path: pathlib.Path, pptx_path: pathlib.Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = parse_slides(md_path.read_text())

    for idx, (title, subtitle, body_lines) in enumerate(slides):
        if idx == 0:
            layout = prs.slide_layouts[TITLE_LAYOUT]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1 and subtitle:
                slide.placeholders[1].text = subtitle
            continue

        layout = prs.slide_layouts[TITLE_CONTENT_LAYOUT]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        build_body(slide, body_lines)

    prs.save(pptx_path)
    print(f"[write] {pptx_path} ({len(slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("usage: md_to_pptx.py <deck.md> <out.pptx>")
        sys.exit(1)
    convert(pathlib.Path(sys.argv[1]), pathlib.Path(sys.argv[2]))
